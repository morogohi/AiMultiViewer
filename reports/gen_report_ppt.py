# -*- coding: utf-8 -*-
"""AI 멀티뷰어 테스트 수행 결과보고서 PPTX 생성 스크립트."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

SHOT_DIR = r"C:\Users\morog\AiMultiViewer\screenshots"
OUT = r"C:\Users\morog\AiMultiViewer\reports\테스트_수행결과보고서_v0.3.0.pptx"

TEAL = RGBColor(0x00, 0x69, 0x6B)
DARK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x5F, 0x6A, 0x6A)
GREEN = RGBColor(0x1E, 0x8E, 0x3E)
AMBER = RGBColor(0xB2, 0x6A, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run()
        r.text = line
        f = r.font
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = color
        f.name = "맑은 고딕"
    return tb


def header(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    add_text(slide, 0.4, 0.12, 11.0, 0.6, title, size=22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    if subtitle:
        add_text(slide, 0.45, 1.02, 12.4, 0.4, subtitle, size=12, color=GRAY)


def shot(slide, name, x=0.5, y=1.5, h=5.75):
    path = os.path.join(SHOT_DIR, name)
    pic = slide.shapes.add_picture(path, Inches(x), Inches(y), height=Inches(h))
    pic.line.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
    pic.line.width = Pt(1)
    return pic


def result_badge(slide, x, y, passed=True, label=None):
    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(1.7), Inches(0.5))
    box.fill.solid()
    box.fill.fore_color.rgb = GREEN if passed else AMBER
    box.line.fill.background()
    tf = box.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = label or ("PASS" if passed else "N/A")
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


# ── 1. 표지 ──────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = TEAL
bg.line.fill.background()
add_text(s, 1.0, 2.3, 11.3, 1.0, "테스트 수행 결과 보고서", size=40, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
add_text(s, 1.0, 3.4, 11.3, 0.6, "AI 멀티뷰어 (Android) v0.3.0 — 문서 판독 · AI 요약 · llm-wiki 자동 수집",
         size=18, color=RGBColor(0xD0, 0xEC, 0xEC), align=PP_ALIGN.CENTER)
add_text(s, 1.0, 5.6, 11.3, 0.9,
         "테스트 일자: 2026-07-12 ~ 07-13   |   환경: Android 에뮬레이터(API 34) · 서명 release APK\n"
         "방법: adb + UI Automator 자동화로 실문서 열람 · 화면 캡처 · 텍스트 덤프 검증",
         size=13, color=RGBColor(0xD0, 0xEC, 0xEC), align=PP_ALIGN.CENTER)

# ── 2. 개요 + 홈 화면 ────────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "1. 테스트 개요", "실제 업무 문서 10종을 앱에서 직접 열람하며 캡처로 검증")
shot(s, "rpt_home.png")
add_text(s, 3.6, 1.6, 9.3, 5.6,
         "■ 테스트 대상\n"
         "  · 실제 업무 문서: PDF 2, DOCX 2, XLSX 2, PPTX 1, HWP 1 (+ ODT/Markdown 샘플)\n"
         "  · 검증 항목: ① 문서 열람(파싱·렌더) ② AI 온디바이스 요약 ③ llm-wiki 자동 수집\n\n"
         "■ 판정 기준\n"
         "  · 열람: 뷰어에 본문 텍스트가 깨짐 없이 표시 (한글 인코딩 포함)\n"
         "  · 수집: 열람 직후 기기 Documents/llm-wiki/에 마크다운 생성 + PC 볼트 동기화\n\n"
         "■ 결과 요약\n"
         "  · 문서 판독 9/9 정상 (HWP는 사양대로 미지원 안내 표시)\n"
         "  · llm-wiki 수집 8건 → Obsidian 볼트 inbox/aimultiviewer/ 반영 확인\n"
         "  · GitHub Actions 릴리스 빌드(v0.3.0) 성공, 서명 APK 배포 완료\n\n"
         "◀ 좌측: 앱 홈 화면 (테스트 문서 목록)", size=15)

# ── 3~12. 문서별 슬라이드 ────────────────────────────────
docs = [
    ("rpt_pdf1.png", "revenuereport.pdf", "PDF", True,
     "월간 수익 리포트 — 표·수치 포함 문서",
     "· 페이지 렌더 + PDFBox 텍스트 추출 정상\n· 수익 요약/충전 실적/설비 현황 표 값 그대로 판독\n"
     "· AI 요약: 핵심 문장 추출 정상 (다음 장 참조)\n· llm-wiki 수집: revenuereport.pdf.md 생성 확인"),
    ("rpt_pdf2.png", "Refund-3693-9270.pdf", "PDF", True,
     "환불/분기 보고 형식의 한글 PDF",
     "· 한글 본문 인코딩 깨짐 없음\n· 문단·불릿 구조 유지\n· llm-wiki 수집: Refund-3693-9270.pdf.md 확인"),
    ("rpt_docx1.png", "251209_SK_Electlink_PoC_Draft.docx", "DOCX", True,
     "영문 제안서 (Hubject Plug&Charge PoC)",
     "· OOXML 파싱 정상 — 표지/목차/본문/표 추출\n"
     "· 개발 중 발견된 숫자 쓰레기(필드코드·도형 좌표) 혼입 문제를\n  XmlTextExtractor 개선으로 수정 완료한 상태로 검증\n"
     "· llm-wiki 수집 확인"),
    ("rpt_docx2.png", "운영지능 온톨로지 기술개발계획서 rev2.docx", "DOCX", True,
     "한글 기술계획서 — 표 다수 포함 (18.7KB 수집)",
     "· 한글 본문·표(클래스 정의, 추론 규칙 등) 정상 추출\n· 장문 문서(약 1만 8천 자) 전문 수집 확인\n"
     "· llm-wiki 수집: ontology_plan_rev2.docx.md 확인"),
    ("rpt_xlsx1.png", "견적서_SK일렉링크_DRMS_구축_20260610.xlsx", "XLSX", True,
     "다중 시트 견적서 (표준견적서/인건비 상세 등 4시트)",
     "· sharedStrings + 시트별 셀 값 추출 정상\n· '── 시트: 표준견적서 ──' 형태로 시트 구분 표시\n"
     "· 금액·수식 결과값 정상 (#REF! 등 원본 오류도 그대로 표기)\n· llm-wiki 수집 확인"),
    ("rpt_xlsx2.png", "에스케이일렉링크 명세서.xlsx", "XLSX", True,
     "거래 명세서",
     "· 시트 셀 값 추출 정상\n· llm-wiki 수집: 에스케이일렉링크_명세서.xlsx.md 확인"),
    ("rpt_pptx.png", "ChargeIQ_개선보고_M6M7.pptx", "PPTX", True,
     "개선보고 프레젠테이션 (슬라이드 다수)",
     "· 슬라이드별 텍스트 추출 — '── 슬라이드 N ──' 구분 표시\n· 제목/불릿/성과 지표 텍스트 정상\n"
     "· 발표자 노트 추출 지원\n· llm-wiki 수집: 9.9KB 마크다운 확인"),
    ("rpt_odt.png", "sample_gdoc.odt", "ODT", True,
     "OpenDocument (Google Docs 내보내기 호환)",
     "· content.xml 텍스트 추출 정상\n· 한글+영문 혼합 텍스트 정상\n"
     "· Google Docs 가상 파일은 동일 경로(PDF 내보내기)로 지원\n  (Drive 계정 필요 — 실기기에서 최종 확인 권장)"),
    ("rpt_md.png", "sample.md", "Markdown", True,
     "마크다운 문서",
     "· 헤딩/불릿 경량 렌더링 정상\n· llm-wiki 수집: sample.md.md 확인"),
    ("rpt_hwp.png", "작업계_볼라드스토퍼_이마트신제주점.hwp", "HWP", False,
     "구형 바이너리 HWP — 미지원 (사양)",
     "· 예상 동작: '미지원 형식' 안내 메시지 표시 → 정상 동작 확인\n"
     "· HWPX(한글 개방형)로 저장 시 열람 가능\n· 네이티브 HWP 파서는 로드맵 M2 예정"),
]
for i, (img, name, fmt, passed, subtitle, detail) in enumerate(docs):
    s = prs.slides.add_slide(BLANK)
    header(s, f"{i + 2}. [{fmt}] {name}", subtitle)
    shot(s, img)
    result_badge(s, 3.6, 1.6, passed, label=("PASS" if passed else "사양(미지원 안내)"))
    add_text(s, 3.6, 2.4, 9.3, 4.6, detail + "\n\n◀ 좌측: 실제 열람 화면 캡처 (에뮬레이터)", size=15)

# ── AI 기능 ──────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "12. AI 기능 — 온디바이스 자동 요약", "revenuereport.pdf 열람 → AI 탭 → 요약 실행")
shot(s, "rpt_ai.png")
result_badge(s, 3.6, 1.5, True)
add_text(s, 3.6, 2.3, 9.3, 4.6,
         "· 네트워크·API 키 없이 온디바이스 추출 요약 즉시 동작\n"
         "· 문서 핵심 수치(충전기별 실적) 중심으로 요약 생성\n"
         "· 클라우드 AI(OpenAI 호환) 활성화 시 고품질 추상 요약·Q&A·비교·회의록·기술분석 확장 가능\n\n"
         "· 동일 요약 엔진이 llm-wiki 수집 마크다운의 '자동 요약' 섹션에도 사용됨\n\n"
         "◀ 좌측: AI 탭 요약 실행 결과 화면", size=15)

# ── llm-wiki 수집 ────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "13. llm-wiki 자동 수집 (v0.3.0 신규)", "열람 → 마크다운 자동 변환 → 기기 축적 → PC Obsidian 볼트 동기화")
shot(s, "rpt_settings.png")
result_badge(s, 3.6, 1.5, True)
add_text(s, 3.6, 2.3, 9.3, 4.8,
         "■ E2E 검증 결과 (전 단계 통과)\n"
         "  ① 문서 열람 즉시 기기 Documents/llm-wiki/에 .md 자동 생성 — 8건 확인\n"
         "  ② 구조 검증: YAML 프런트매터 + 자동 요약 + 전문, 한글 UTF-8 정상\n"
         "  ③ 재열람 시 동일 파일 갱신 → 중복 없음\n"
         "  ④ sync_wiki_from_device.ps1 → 볼트 inbox/aimultiviewer/ 8건 수집 (한글 파일명 보존)\n"
         "  ⑤ 작업 스케줄러 'llm-wiki-sync' 등록 — 매시간 자동 동기화\n"
         "  ⑥ 설정 '열람 문서 자동 수집' 토글 on/off 동작\n\n"
         "■ 수집 현황: 견적서/명세서(XLSX), 제안서/계획서(DOCX), 개선보고(PPTX),\n"
         "   수익리포트/환불(PDF), 샘플(MD) — 총 8건, 최대 18.7KB\n\n"
         "◀ 좌측: 설정 화면 — llm-wiki 자동 수집 토글(켜짐)", size=14)

# ── 종합 결과 ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "14. 종합 결과", "v0.3.0 기준 — 2026-07-13")
rows = [
    ("구분", "항목", "결과"),
    ("문서 판독", "PDF·DOCX·XLSX·PPTX·ODT·Markdown 9종 열람", "PASS (9/9)"),
    ("문서 판독", "HWP 구형 바이너리 미지원 안내", "PASS (사양)"),
    ("AI", "온디바이스 자동 요약", "PASS"),
    ("llm-wiki", "자동 변환·축적 + 볼트 동기화 E2E 6단계", "PASS (6/6)"),
    ("배포", "GitHub Actions v0.3.0 서명 릴리스", "PASS"),
]
tbl = s.shapes.add_table(len(rows), 3, Inches(0.7), Inches(1.4), Inches(12.0), Inches(3.2)).table
tbl.columns[0].width = Inches(2.2)
tbl.columns[1].width = Inches(7.3)
tbl.columns[2].width = Inches(2.5)
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(14)
        p.runs[0].font.name = "맑은 고딕"
        p.runs[0].font.bold = (r == 0)
        if r > 0 and c == 2:
            p.runs[0].font.color.rgb = GREEN
            p.runs[0].font.bold = True
add_text(s, 0.7, 5.1, 12.0, 1.8,
         "결론: 계획된 기능이 모두 정상 동작함을 실문서 기반으로 확인. 잔여 과제는 HWP/DOC 네이티브 파서(M2),\n"
         "온디바이스 LLM 기반 RAG(M3), Google 문서 실기기 E2E 확인.\n"
         "※ 본 보고서의 캡처에는 실제 업무 문서 내용이 포함되어 있어 공개 저장소(GitHub)에는 게시하지 않음 (로컬 보관).",
         size=14, color=GRAY)

prs.save(OUT)
print("SAVED:", OUT)
