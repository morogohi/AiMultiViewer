# -*- coding: utf-8 -*-
"""AI 멀티뷰어 v0.4.0 문서 보기 기능 재검증 결과보고서 PPTX 생성."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

SHOT_DIR = r"C:\Users\morog\AiMultiViewer\screenshots"
OUT = r"C:\Users\morog\AiMultiViewer\reports\테스트_수행결과보고서_v0.4.0.pptx"

TEAL = RGBColor(0x00, 0x69, 0x6B)
DARK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x5F, 0x6A, 0x6A)
GREEN = RGBColor(0x1E, 0x8E, 0x3E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run()
        r.text = line
        f = r.font
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = color
        f.name = "맑은 고딕"
    return tb


def header(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    add_text(slide, 0.4, 0.12, 12.4, 0.6, title, size=22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    if subtitle:
        add_text(slide, 0.45, 1.02, 12.4, 0.4, subtitle, size=12, color=GRAY)


def shot(slide, name, x=0.5, y=1.5, h=5.75):
    pic = slide.shapes.add_picture(os.path.join(SHOT_DIR, name), Inches(x), Inches(y), height=Inches(h))
    pic.line.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
    pic.line.width = Pt(1)
    return pic


def badge(slide, x, y, label="PASS", color=GREEN):
    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(1.7), Inches(0.5))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = label
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


# ── 표지 ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = TEAL
bg.line.fill.background()
add_text(s, 1.0, 2.2, 11.3, 1.0, "문서 보기 기능 재검증 결과 보고서", size=40, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
add_text(s, 1.0, 3.3, 11.3, 0.6,
         "AI 멀티뷰어 (Android) v0.4.0 — 구조화 렌더링(표 그리드) · HWP 5.0 네이티브 파서",
         size=18, color=RGBColor(0xD0, 0xEC, 0xEC), align=PP_ALIGN.CENTER)
add_text(s, 1.0, 5.6, 11.3, 0.9,
         "테스트 일자: 2026-07-13   |   환경: Android 에뮬레이터(API 34) · 서명 release APK (v0.4.0)\n"
         "방법: 전체 문서 11종을 개별 열람하며 실화면 캡처로 검증 (adb + UI Automator 자동화)",
         size=13, color=RGBColor(0xD0, 0xEC, 0xEC), align=PP_ALIGN.CENTER)

# ── 개요 ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "1. 테스트 개요", "v0.4.0 개선 후 전체 포맷 재검증 — 문서 11종 전수 열람")
shot(s, "t_home.png")
add_text(s, 3.6, 1.6, 9.3, 5.6,
         "■ v0.4.0 개선 사항 (이번 검증 대상)\n"
         "  ① PPT·워드·엑셀 구조화 렌더링 — 제목·불릿·표를 구조 그대로 표시\n"
         "      (엑셀·문서 내 표 → 테두리·헤더행·가로 스크롤을 갖춘 표 그리드)\n"
         "  ② 구형 HWP(5.0 바이너리) 자체 파서 신규 — HWPX 변환 없이 바로 열람\n\n"
         "■ 테스트 범위 — 등록 문서 11종 전수\n"
         "  · PDF 2, DOCX 2, XLSX 2, PPTX 1, HWP 1, ODT 1, 이미지 1, Markdown 1\n\n"
         "■ 판정 기준\n"
         "  · 열람 성공 + 구조 보존(제목/표/불릿) + 한글 인코딩 정상\n\n"
         "■ 결과: 11/11 PASS (전 문서 정상 열람, 실화면 캡처 첨부)\n\n"
         "◀ 좌측: 앱 홈 — 테스트 문서 목록", size=15)

# ── 문서별 ───────────────────────────────────────────────
docs = [
    ("t_xlsx1.png", "[XLSX] 견적서_SK일렉링크_DRMS_구축", "엑셀 → 표 그리드 렌더링 (v0.4.0 핵심 개선)",
     "· '시트: 표준견적서' 제목 아래 셀이 실제 표 그리드로 표시\n"
     "· 테두리·헤더행 강조·열 정렬, 셀 위치(r=\"B3\") 기반 열 배치\n"
     "· 넓은 표는 가로 스크롤로 열람\n"
     "· 이전(v0.3): 탭 구분 텍스트 나열 → 현재: 표 구조 그대로"),
    ("t_xlsx2.png", "[XLSX] 에스케이일렉링크_명세서", "거래명세서 — 병합 셀 포함 양식 문서",
     "· 거래명세서 양식(등록번호/상호/대표 등)이 표 형태로 정돈되어 표시\n"
     "· 합계금액·품목 행 정상\n"
     "· 빈 셀은 공백 셀로 유지되어 열 정렬 보존"),
    ("t_docx1.png", "[DOCX] 251209_SK_Electlink_PoC_Draft", "워드 → 문단 + 문서 내 표 그리드",
     "· 본문 문단과 문서 내 표(Service Provider/Category/Contact…)가\n"
     "  헤더행이 강조된 표 그리드로 렌더링\n"
     "· 셀 안 여러 줄(주소·이메일)도 셀 내 줄바꿈으로 보존\n"
     "· 필드코드·도형 좌표 등 비표시 데이터 혼입 없음"),
    ("t_docx2.png", "[DOCX] 운영지능 온톨로지 기술개발계획서", "한글 장문 문서 — 문단 + 비교표",
     "· 한글 본문 문단 정상, '관계형 DB vs 온톨로지' 비교표가\n"
     "  표 그리드로 표시 (구분/데이터 표현/지식 활용/장애 대응…)\n"
     "· 약 1만 8천 자 장문도 스크롤 열람 원활"),
    ("t_pptx.png", "[PPTX] ChargeIQ_개선보고_M6M7", "파워포인트 → 슬라이드 카드 + 불릿",
     "· '슬라이드 N' 제목으로 슬라이드 경계 구분\n"
     "· 본문 텍스트는 불릿 목록으로 정돈\n"
     "· 발표자 노트는 회색 박스로 첨부 (해당 슬라이드 하단)"),
    ("t_hwp.png", "[HWP] 작업계_볼라드스토퍼_이마트신제주점", "구형 HWP 5.0 — v0.4.0 신규 지원",
     "· 이전: '미지원 형식' 안내 → 현재: 본문 전체 열람\n"
     "· 자체 CFB 리더 + 레코드 파서 (외부 라이브러리 없음)\n"
     "· 작업 시간·작업 내용·작업자 연락처 등 표 셀 텍스트까지 추출\n"
     "· 암호화/배포용(DRM) 문서는 명확한 안내 표시"),
    ("t_pdf1.png", "[PDF] revenuereport", "PDF 페이지 렌더 — 표·수치 문서",
     "· PdfRenderer 페이지 렌더 정상 (표·서식 원본 그대로)\n"
     "· 텍스트 추출(PDFBox) 병행 → AI 요약·위키 수집에 사용"),
    ("t_pdf2.png", "[PDF] Refund-3693-9270", "영수증형 PDF",
     "· 페이지 렌더 정상, 결제 내역·환불 정보 원본 레이아웃 유지"),
    ("t_odt.png", "[ODT] sample_gdoc", "OpenDocument (Google Docs 내보내기 호환)",
     "· content.xml 텍스트 추출 정상, 한·영 혼용 정상"),
    ("t_img.png", "[이미지] 테스트_스크린샷.png", "이미지 표시 + 온디바이스 OCR",
     "· 원본 이미지 표시 정상\n"
     "· ML Kit OCR(한국어+영문)로 이미지 속 텍스트 추출 → AI·위키 활용"),
    ("t_md.png", "[Markdown] sample.md", "마크다운 경량 렌더링",
     "· 헤딩 크기 구분·불릿 렌더링 정상"),
]
for i, (img, title, subtitle, detail) in enumerate(docs):
    s = prs.slides.add_slide(BLANK)
    header(s, f"{i + 2}. {title}", subtitle)
    shot(s, img)
    badge(s, 3.6, 1.6)
    add_text(s, 3.6, 2.4, 9.3, 4.6, detail + "\n\n◀ 좌측: 실제 열람 화면 캡처 (v0.4.0, 에뮬레이터)", size=15)

# ── 종합 ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "13. 종합 결과", "v0.4.0 · 2026-07-13 · 문서 11종 전수 재검증")
rows = [
    ("포맷", "문서", "검증 포인트", "결과"),
    ("XLSX", "견적서 DRMS / 명세서", "표 그리드·열 정렬·가로 스크롤", "PASS"),
    ("DOCX", "PoC Draft / 온톨로지 계획서", "문단 + 문서 내 표 그리드, 쓰레기 혼입 없음", "PASS"),
    ("PPTX", "ChargeIQ 개선보고", "슬라이드 구분 + 불릿 + 노트", "PASS"),
    ("HWP", "작업계 (구형 5.0)", "신규 자체 파서로 본문 열람", "PASS"),
    ("PDF", "수익리포트 / 환불영수증", "페이지 렌더 + 텍스트 추출", "PASS"),
    ("ODT/이미지/MD", "샘플 3종", "텍스트/OCR/헤딩 렌더", "PASS"),
]
tbl = s.shapes.add_table(len(rows), 4, Inches(0.7), Inches(1.5), Inches(12.0), Inches(3.4)).table
tbl.columns[0].width = Inches(1.9)
tbl.columns[1].width = Inches(3.6)
tbl.columns[2].width = Inches(4.8)
tbl.columns[3].width = Inches(1.7)
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(13)
        p.runs[0].font.name = "맑은 고딕"
        p.runs[0].font.bold = (r == 0)
        if r > 0 and c == 3:
            p.runs[0].font.color.rgb = GREEN
            p.runs[0].font.bold = True
add_text(s, 0.7, 5.3, 12.0, 1.7,
         "결론: v0.4.0 개선(구조화 렌더링·HWP 지원) 후 등록 문서 11종 전수 열람 재검증 — 전 항목 PASS.\n"
         "잔여 과제: DOC(Word 97) 파서, HWP 표 구조(그리드) 렌더링, Google 문서 실기기 E2E.\n"
         "※ 캡처에 실제 업무 문서 내용이 포함되어 본 보고서는 로컬 보관 (공개 저장소 미게시).",
         size=14, color=GRAY)

prs.save(OUT)
print("SAVED:", OUT, "slides:", len(prs.slides._sldIdLst))
